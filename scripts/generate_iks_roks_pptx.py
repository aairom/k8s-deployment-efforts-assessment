#!/usr/bin/env python3
"""
Génère la présentation comparative IKS vs ROKS vs Self-Provisioned Kubernetes vs Amazon EKS.
Sauvegardée sous : IKS-vs-ROKS-Analyse-Comparative.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ── Palette de couleurs ──────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1E, 0x27, 0x61)   # IBM bleu foncé
MED_BLUE    = RGBColor(0x00, 0x61, 0xAA)   # IBM bleu moyen
LIGHT_BLUE  = RGBColor(0xD0, 0xE4, 0xF5)   # fond bleu clair
RED_OCP     = RGBColor(0xCC, 0x00, 0x00)   # Rouge OpenShift
ORANGE_EKS  = RGBColor(0xFF, 0x99, 0x00)   # Orange AWS EKS
GREEN_K8S   = RGBColor(0x32, 0x6C, 0xE5)   # Bleu Kubernetes
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY    = RGBColor(0x76, 0x76, 0x76)
ACCENT_GOLD = RGBColor(0xF5, 0xA6, 0x23)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        if line_width:
            line.width = Pt(line_width)
    else:
        line.fill.background()
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_size=18, bold=False, italic=False,
                 color=DARK_GRAY, align=PP_ALIGN.LEFT,
                 font_name="Calibri", word_wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_text_box_multiline(slide, lines, x, y, w, h,
                           font_size=16, bold=False, color=DARK_GRAY,
                           font_name="Calibri", line_spacing=None, bullet=False):
    """lines: list of str or (str, bold, size, color) tuples"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, is_bold, fs, col = item, bold, font_size, color
        else:
            text = item[0]
            is_bold = item[1] if len(item) > 1 else bold
            fs = item[2] if len(item) > 2 else font_size
            col = item[3] if len(item) > 3 else color

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.alignment = PP_ALIGN.LEFT
        if line_spacing:
            p.line_spacing = Pt(line_spacing)

        run = p.add_run()
        run.text = text
        run.font.size = Pt(fs)
        run.font.bold = is_bold
        run.font.color.rgb = col
        run.font.name = font_name

    return txBox


def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_slide_header_bar(slide, title, bg_color=DARK_BLUE, title_color=WHITE,
                         bar_h=1.1, font_size=30):
    add_rect(slide, 0, 0, 13.33, bar_h, fill_color=bg_color)
    add_text_box(slide, title, 0.3, 0.08, 12.5, bar_h - 0.1,
                 font_size=font_size, bold=True, color=title_color,
                 align=PP_ALIGN.LEFT, font_name="Calibri")


def add_footer(slide, text="IKS | ROKS | Kubernetes | EKS — Analyse Comparative © 2026",
               page_num=None):
    # Footer bar
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill_color=DARK_BLUE)
    label = text
    if page_num:
        label = f"{text}    |    {page_num}"
    add_text_box(slide, label, 0.2, 7.22, 12.9, 0.28,
                 font_size=9, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")


def add_colored_bullet(slide, items, x, y, w, h,
                       font_size=16, bullet_color=MED_BLUE, text_color=DARK_GRAY,
                       font_name="Calibri", title=None, title_color=DARK_BLUE):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True

    if title:
        p = tf.paragraphs[0]
        first = False
        run = p.add_run()
        run.text = title
        run.font.size = Pt(font_size + 2)
        run.font.bold = True
        run.font.color.rgb = title_color
        run.font.name = font_name

    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = text_color
        run.font.name = font_name
    return txBox


def add_table(slide, data, x, y, w, h,
              header_fill=DARK_BLUE, header_text=WHITE,
              row_alt=LIGHT_BLUE, font_size=13):
    """data: list of rows (list of str). First row = header."""
    rows = len(data)
    cols = len(data[0])
    tbl = slide.shapes.add_table(rows, cols,
                                 Inches(x), Inches(y),
                                 Inches(w), Inches(h)).table

    col_width = Inches(w / cols)
    for c in range(cols):
        tbl.columns[c].width = col_width

    for r, row in enumerate(data):
        for c, cell_text in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(cell_text)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(font_size)
            run.font.name = "Calibri"

            if r == 0:
                run.font.bold = True
                run.font.color.rgb = header_text
                fill = cell.fill
                fill.solid()
                fill.fore_color.rgb = header_fill
            else:
                run.font.color.rgb = DARK_GRAY
                if r % 2 == 0:
                    fill = cell.fill
                    fill.solid()
                    fill.fore_color.rgb = row_alt
                else:
                    fill = cell.fill
                    fill.solid()
                    fill.fore_color.rgb = WHITE
    return tbl


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Diapositive 1 — Titre"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide, DARK_BLUE)

    # Bande décorative
    add_rect(slide, 0, 5.8, 13.33, 0.08, fill_color=ACCENT_GOLD)

    # Titre principal
    add_text_box(slide,
                 "Analyse Comparative des Solutions Kubernetes",
                 0.6, 1.0, 12.0, 1.6,
                 font_size=38, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font_name="Calibri")

    # Sous-titre
    add_text_box(slide,
                 "IKS · ROKS · Kubernetes Auto-Géré · Amazon EKS",
                 0.6, 2.7, 12.0, 0.8,
                 font_size=24, bold=False, color=ACCENT_GOLD,
                 align=PP_ALIGN.CENTER, font_name="Calibri")

    # Description
    add_text_box(slide,
                 "Évaluation multicritère pour la sélection d'une plateforme de conteneurs",
                 0.6, 3.55, 12.0, 0.6,
                 font_size=16, bold=False, color=LIGHT_BLUE,
                 align=PP_ALIGN.CENTER, font_name="Calibri")

    # Date
    add_text_box(slide,
                 "3 septembre 2026",
                 0.6, 4.4, 12.0, 0.5,
                 font_size=14, bold=False, color=MID_GRAY,
                 align=PP_ALIGN.CENTER, font_name="Calibri")

    # Logos colorés des solutions
    colors = [MED_BLUE, RED_OCP, GREEN_K8S, ORANGE_EKS]
    labels = ["IKS", "ROKS", "K8s Auto-Géré", "Amazon EKS"]
    for i, (col, lbl) in enumerate(zip(colors, labels)):
        bx = 1.5 + i * 2.8
        add_rect(slide, bx, 5.1, 2.2, 0.55, fill_color=col)
        add_text_box(slide, lbl, bx, 5.1, 2.2, 0.55,
                     font_size=14, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, font_name="Calibri")

    add_footer(slide)


def slide_02_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_slide_header_bar(slide, "Sommaire")

    sections = [
        ("01", "Introduction & Contexte",           "Pourquoi comparer ces solutions ?"),
        ("02", "Vue d'ensemble des Solutions",       "IKS · ROKS · K8s Auto-Géré · EKS"),
        ("03", "Analyse Comparative — Dimensions",   "Infrastructure, Déploiement, Sécurité, Réseau…"),
        ("04", "Tableau Récapitulatif",              "Comparaison synthétique côte à côte"),
        ("05", "Recommandations",                    "Quel outil pour quel profil ?"),
        ("06", "Conclusion",                         "Points clés à retenir"),
    ]

    for i, (num, title, subtitle) in enumerate(sections):
        y = 1.25 + i * 0.97
        add_rect(slide, 0.4, y, 0.65, 0.72, fill_color=DARK_BLUE)
        add_text_box(slide, num, 0.4, y, 0.65, 0.72,
                     font_size=20, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, font_name="Calibri")
        add_text_box(slide, title, 1.2, y + 0.02, 6.5, 0.38,
                     font_size=17, bold=True, color=DARK_BLUE, font_name="Calibri")
        add_text_box(slide, subtitle, 1.2, y + 0.38, 6.5, 0.3,
                     font_size=13, bold=False, color=MID_GRAY, font_name="Calibri")

    add_footer(slide, page_num="2")


def slide_03_intro(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_slide_header_bar(slide, "01 — Introduction & Contexte")

    add_text_box(slide,
                 "Pourquoi évaluer les plateformes Kubernetes managées et auto-gérées ?",
                 0.4, 1.25, 12.5, 0.55,
                 font_size=19, bold=True, color=DARK_BLUE, font_name="Calibri")

    bullets_left = [
        "La conteneurisation est devenue le standard pour le déploiement cloud-native.",
        "Les organisations cherchent à réduire la charge opérationnelle liée à la gestion de Kubernetes.",
        "Les exigences réglementaires, de souveraineté des données et de conformité influencent le choix.",
        "Les coûts, la flexibilité et l'intégration écosystémique sont des facteurs décisifs.",
        "La stratégie multi-cloud ou hybride nécessite une évaluation objective des solutions disponibles.",
    ]
    add_colored_bullet(slide, bullets_left, 0.4, 1.95, 6.2, 4.9,
                       font_size=15, font_name="Calibri")

    # Boîte de contexte à droite
    add_rect(slide, 6.9, 1.95, 6.0, 4.9, fill_color=LIGHT_BLUE)
    add_text_box(slide, "Périmètre de cette analyse", 7.1, 2.1, 5.7, 0.45,
                 font_size=16, bold=True, color=DARK_BLUE, font_name="Calibri")
    context = [
        "IBM Kubernetes Service (IKS)\n  → Kubernetes managé sur IBM Cloud",
        "Red Hat OpenShift on IBM Cloud (ROKS)\n  → OpenShift managé sur IBM Cloud",
        "Kubernetes Auto-Géré (On-Prem / Self-Managed)\n  → Infrastructure propre ou cloud non managé",
        "Amazon EKS\n  → Kubernetes managé sur AWS",
    ]
    for j, ctx in enumerate(context):
        y_ctx = 2.65 + j * 0.95
        add_rect(slide, 7.05, y_ctx, 5.6, 0.78, fill_color=WHITE)
        add_text_box(slide, ctx, 7.15, y_ctx + 0.02, 5.4, 0.74,
                     font_size=13, bold=False, color=DARK_GRAY, font_name="Calibri")

    add_footer(slide, page_num="3")


def _solution_slide(prs, title, bar_color, provider, technology, version_k8s,
                    description, characteristics, use_cases, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_slide_header_bar(slide, title, bg_color=bar_color)

    # Description
    add_text_box(slide, description, 0.4, 1.2, 12.5, 0.65,
                 font_size=15, bold=False, color=DARK_GRAY, font_name="Calibri")

    # Infos rapides
    quick = [
        ("Fournisseur",          provider),
        ("Technologie de base",  technology),
        ("Version Kubernetes",   version_k8s),
    ]
    for i, (label, val) in enumerate(quick):
        bx = 0.4 + i * 4.3
        add_rect(slide, bx, 1.95, 4.0, 0.7, fill_color=bar_color)
        add_text_box(slide, label, bx, 1.95, 4.0, 0.35,
                     font_size=11, bold=False, color=WHITE,
                     align=PP_ALIGN.CENTER, font_name="Calibri")
        add_text_box(slide, val, bx, 2.3, 4.0, 0.35,
                     font_size=13, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, font_name="Calibri")

    # Caractéristiques
    add_text_box(slide, "Caractéristiques clés", 0.4, 2.85, 6.2, 0.4,
                 font_size=16, bold=True, color=bar_color, font_name="Calibri")
    add_colored_bullet(slide, characteristics, 0.4, 3.3, 6.2, 3.5,
                       font_size=14, bullet_color=bar_color, font_name="Calibri")

    # Cas d'usage
    add_text_box(slide, "Cas d'usage typiques", 7.0, 2.85, 5.9, 0.4,
                 font_size=16, bold=True, color=bar_color, font_name="Calibri")
    add_colored_bullet(slide, use_cases, 7.0, 3.3, 5.9, 3.5,
                       font_size=14, bullet_color=bar_color, font_name="Calibri")

    add_footer(slide, page_num=str(page_num))


def slide_04_iks(prs):
    _solution_slide(
        prs,
        title="02 — IBM Kubernetes Service (IKS)",
        bar_color=MED_BLUE,
        provider="IBM Cloud",
        technology="Kubernetes vanilla (CNCF certifié)",
        version_k8s="1.29+ (mises à jour automatiques disponibles)",
        description=(
            "IKS est le service Kubernetes entièrement managé d'IBM Cloud, offrant un plan de contrôle "
            "géré par IBM, avec des nœuds worker déployés dans l'infrastructure IBM Cloud."
        ),
        characteristics=[
            "Plan de contrôle Kubernetes entièrement managé par IBM",
            "Intégration native avec IBM Cloud IAM, IBM Cloud Object Storage",
            "Support des zones de disponibilité multiples (multi-zone)",
            "Intégration avec IBM Cloud Monitoring, Logging et Security Advisor",
            "Mise à jour automatique ou à la demande du plan de contrôle",
            "Accès complet à l'API Kubernetes standard",
            "Support des GPU et architectures spécialisées",
            "Conformité ISO 27001, SOC 2, PCI DSS, HIPAA",
        ],
        use_cases=[
            "Applications cloud-native sur IBM Cloud",
            "Organisations préférant Kubernetes standard sans couche OpenShift",
            "Workloads nécessitant une intégration IBM Cloud (Watson, Db2, etc.)",
            "Migration de workloads on-premises vers IBM Cloud",
            "Applications nécessitant haute disponibilité multi-zone",
            "Équipes DevOps recherchant un accès kubectl natif",
        ],
        page_num=4
    )


def slide_05_roks(prs):
    _solution_slide(
        prs,
        title="02 — Red Hat OpenShift on IBM Cloud (ROKS)",
        bar_color=RED_OCP,
        provider="IBM Cloud (avec Red Hat OpenShift)",
        technology="OpenShift Container Platform (OCP) sur Kubernetes",
        version_k8s="OCP 4.x (basé sur K8s 1.28+)",
        description=(
            "ROKS combine la puissance d'OpenShift Container Platform de Red Hat avec la fiabilité "
            "et les services d'IBM Cloud, dans un environnement entièrement managé."
        ),
        characteristics=[
            "OpenShift entièrement managé — plan de contrôle pris en charge par IBM",
            "Outils développeur intégrés : oc CLI, Developer Console, Operator Hub",
            "Pipelines CI/CD natifs (OpenShift Pipelines / Tekton, ArgoCD)",
            "Sécurité renforcée : SCC (Security Context Constraints) par défaut",
            "Intégration Red Hat Marketplace et Operator Catalog",
            "Registry d'images intégré (OpenShift Internal Registry)",
            "Routes OpenShift + Ingress Kubernetes supportés",
            "Conformité FIPS 140-2, FedRAMP, DoD CC SRG (option GovCloud)",
        ],
        use_cases=[
            "Entreprises utilisant déjà l'écosystème Red Hat (RHEL, Ansible)",
            "Équipes cherchant une expérience développeur enrichie via OpenShift",
            "Applications nécessitant conformité FedRAMP ou DoD",
            "Projets nécessitant pipelines CI/CD intégrés dès le départ",
            "Workloads hybrides couplant OpenShift on-premises et cloud",
            "Migration de workloads Red Hat Enterprise Linux vers le cloud",
        ],
        page_num=5
    )


def slide_06_k8s_self(prs):
    _solution_slide(
        prs,
        title="02 — Kubernetes Auto-Géré (On-Premises / Self-Managed)",
        bar_color=GREEN_K8S,
        provider="Auto-géré (Bare Metal, VMware, OpenStack, etc.)",
        technology="Kubernetes vanilla ou distribution (kubeadm, k3s, RKE, Rancher…)",
        version_k8s="Choisie par l'organisation (contrôle total)",
        description=(
            "Kubernetes auto-géré désigne tout déploiement où l'organisation prend en charge "
            "l'installation, la configuration, les mises à jour et l'exploitation de Kubernetes, "
            "sur une infrastructure propre ou non managée."
        ),
        characteristics=[
            "Contrôle total sur la version, la configuration et les composants",
            "Déploiement sur infrastructure propre, VMware, OpenStack ou cloud IaaS",
            "Flexibilité maximale sur le choix du CNI, CSI, distribution",
            "Aucune dépendance fournisseur (vendor lock-in minimal)",
            "Charge opérationnelle totale à la charge de l'équipe interne",
            "Intégration possible avec outils existants (Ansible, Terraform, GitOps)",
            "Peut être couplé à des outils enterprise (Rancher, Tanzu, OpenShift autogéré)",
            "Certification CNCF possible selon la distribution choisie",
        ],
        use_cases=[
            "Organisations avec exigences de souveraineté des données strictes",
            "Environnements air-gapped (déconnectés d'Internet)",
            "Datacenter on-premises existant à rentabiliser",
            "Équipes SRE/Ops expérimentées souhaitant un contrôle total",
            "Workloads réglementés ne pouvant pas aller dans le cloud public",
            "Prototypage ou apprentissage de Kubernetes",
        ],
        page_num=6
    )


def slide_07_eks(prs):
    _solution_slide(
        prs,
        title="02 — Amazon Elastic Kubernetes Service (EKS)",
        bar_color=ORANGE_EKS,
        provider="Amazon Web Services (AWS)",
        technology="Kubernetes managé AWS (CNCF certifié)",
        version_k8s="1.28+ (mises à jour standard AWS)",
        description=(
            "Amazon EKS est le service Kubernetes managé d'AWS, permettant de déployer et "
            "d'exploiter des clusters Kubernetes sur l'infrastructure AWS sans gérer le plan de contrôle."
        ),
        characteristics=[
            "Plan de contrôle multi-AZ entièrement managé par AWS",
            "Intégration native avec IAM AWS, VPC, ALB, EBS, EFS, ECR",
            "Support des nœuds EC2 (managed node groups) et AWS Fargate (serverless)",
            "EKS Anywhere pour déploiement on-premises",
            "Add-ons EKS gérés : CoreDNS, kube-proxy, VPC CNI, etc.",
            "Intégration AWS GuardDuty, AWS Security Hub, Macie",
            "Support ARM (Graviton) et GPU (NVIDIA, Inferentia)",
            "Conformité SOC, PCI, ISO, HIPAA, FedRAMP",
        ],
        use_cases=[
            "Organisations natives AWS utilisant intensivement les services AWS",
            "Applications nécessitant intégration S3, RDS, DynamoDB, Lambda",
            "Workloads serverless via EKS + Fargate",
            "Équipes cherchant un large écosystème de solutions partenaires AWS",
            "Déploiements multi-régions sur AWS avec Transit Gateway",
            "Migrations depuis ECS vers Kubernetes sur AWS",
        ],
        page_num=7
    )


def _comparison_slide(prs, section_num, section_title, dimension,
                      iks_text, roks_text, k8s_text, eks_text, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_slide_header_bar(slide, f"{section_num} — Comparaison : {dimension}")

    # 4 colonnes
    configs = [
        ("IKS",            MED_BLUE,   iks_text),
        ("ROKS",           RED_OCP,    roks_text),
        ("K8s Auto-Géré",  GREEN_K8S,  k8s_text),
        ("Amazon EKS",     ORANGE_EKS, eks_text),
    ]

    col_w = 3.1
    for i, (label, col, content) in enumerate(configs):
        bx = 0.2 + i * 3.3
        add_rect(slide, bx, 1.2, col_w, 0.5, fill_color=col)
        add_text_box(slide, label, bx, 1.2, col_w, 0.5,
                     font_size=16, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, font_name="Calibri")
        add_rect(slide, bx, 1.7, col_w, 5.3, fill_color=LIGHT_GRAY)
        add_text_box_multiline(slide, content, bx + 0.1, 1.8, col_w - 0.2, 5.1,
                               font_size=13, color=DARK_GRAY, font_name="Calibri")

    add_footer(slide, page_num=str(page_num))


def slides_comparative(prs):
    comparisons = [
        {
            "section_num": "03",
            "section_title": "Infrastructure",
            "dimension": "Infrastructure & Modèle d'hébergement",
            "iks": [
                "• IBM Cloud (IaaS virtuel)",
                "• Nœuds VMs sur datacenters IBM",
                "• Multi-zone dans une région IBM Cloud",
                "• Bare metal disponible en option",
                "• Isolation VPC ou Classic Network",
                "• SLA 99,9% sur plan de contrôle",
            ],
            "roks": [
                "• IBM Cloud (IaaS virtuel)",
                "• Nœuds VMs ou bare metal IBM Cloud",
                "• Multi-zone avec option ROSA (AWS)",
                "• Architecture OCP 4.x avec etcd géré",
                "• Isolation VPC ou Classic Network",
                "• SLA 99,9% plan de contrôle",
            ],
            "k8s": [
                "• Infrastructure propre (on-prem, VMware, bare metal)",
                "• Contrôle total matériel et réseau",
                "• Topologie entièrement personnalisable",
                "• Aucune limite géographique fournisseur",
                "• Responsabilité SLA interne",
                "• Haute disponibilité à concevoir soi-même",
            ],
            "eks": [
                "• AWS (EC2, Fargate, Outposts)",
                "• Managed Node Groups ou nœuds autogérés",
                "• Multi-AZ dans régions AWS (35+ régions)",
                "• EKS Anywhere pour on-premises",
                "• Intégration VPC native AWS",
                "• SLA 99,95% plan de contrôle",
            ],
        },
        {
            "section_num": "03",
            "section_title": "Déploiement",
            "dimension": "Facilité de déploiement & Provisionnement",
            "iks": [
                "• CLI ibmcloud, Terraform, UI IBM Cloud",
                "• Cluster opérationnel en ~10-15 min",
                "• Auto-scaling nœuds (Cluster Autoscaler)",
                "• Node pools avec profils standardisés",
                "• Intégration Schematics (Terraform managé)",
                "• Blueprints pour architectures prédéfinies",
            ],
            "roks": [
                "• CLI ibmcloud, oc, Terraform, UI IBM Cloud",
                "• Cluster opérationnel en ~20-30 min (OCP plus lourd)",
                "• Machine Config Operator pour OS des nœuds",
                "• OperatorHub pour add-ons en un clic",
                "• Intégration Schematics / Ansible",
                "• Developer Hub intégré (OpenShift DevSpaces)",
            ],
            "k8s": [
                "• kubeadm, k3s, kubespray, RKE, Rancher, Talos…",
                "• Déploiement de plusieurs heures à jours",
                "• Provisionnement manuel ou via Terraform/Ansible",
                "• Gestion des certificats, etcd, PKI manuelle",
                "• Courbe d'apprentissage élevée",
                "• Flexibilité maximale mais complexité accrue",
            ],
            "eks": [
                "• eksctl, AWS Console, Terraform, CDK",
                "• Cluster opérationnel en ~10-15 min",
                "• Managed Node Groups avec mise à jour automatique",
                "• Fargate profiles pour pods serverless",
                "• AWS CloudFormation / CDK pour IaC",
                "• Blueprints EKS (eks-blueprints-cdk)",
            ],
        },
        {
            "section_num": "03",
            "section_title": "Versions Kubernetes",
            "dimension": "Gestion des versions Kubernetes & Mises à jour",
            "iks": [
                "• Support N-2 versions Kubernetes",
                "• Mise à jour manuelle déclenchée par l'utilisateur",
                "• Mise à jour automatique configurable",
                "• Patch de sécurité plan de contrôle par IBM",
                "• Fenêtres de maintenance planifiables",
                "• Versions EOL retirées avec préavis 6 mois",
            ],
            "roks": [
                "• Cycle de version OCP (tous les 6 mois env.)",
                "• Mises à jour OCP managées par IBM",
                "• OpenShift Update Service (OSUS) intégré",
                "• Support Long Term Support (LTS) OCP disponible",
                "• OS nœuds (RHCOS) mis à jour automatiquement",
                "• Validation pré-upgrade automatique",
            ],
            "k8s": [
                "• Version au choix de l'organisation",
                "• Mises à jour entièrement manuelles",
                "• Risque de dérive de version entre nœuds",
                "• Support EOL CNCF : N-3 versions maintenues",
                "• Patching OS nœuds sous responsabilité interne",
                "• Flexibilité : rester sur une ancienne version si besoin",
            ],
            "eks": [
                "• Support N-4 versions Kubernetes (politique étendue)",
                "• Mise à jour plan de contrôle en 1 clic ou automatique",
                "• Extended Support (payant) pour versions plus anciennes",
                "• Patch automatique des Add-ons EKS gérés",
                "• Mise à jour nœuds via Managed Node Group avec drain",
                "• Bottlerocket OS pour nœuds avec mises à jour atomiques",
            ],
        },
        {
            "section_num": "03",
            "section_title": "Sécurité",
            "dimension": "Sécurité & Conformité",
            "iks": [
                "• IAM IBM Cloud + RBAC Kubernetes natif",
                "• Chiffrement etcd au repos (AES-256)",
                "• IBM Key Protect / Hyper Protect Crypto Services",
                "• Scan d'images via IBM Vulnerability Advisor",
                "• Network Policies (Calico) intégrées",
                "• Conformité : ISO 27001, SOC 2, PCI DSS, HIPAA",
            ],
            "roks": [
                "• Security Context Constraints (SCC) OpenShift",
                "• IAM IBM Cloud + OAuth OpenShift",
                "• Pods s'exécutant sans root par défaut",
                "• OPA/Gatekeeper et Kyverno supportés",
                "• Intégration Red Hat ACS (Advanced Cluster Security)",
                "• Conformité FIPS 140-2, FedRAMP, DoD CC SRG",
            ],
            "k8s": [
                "• Sécurité entièrement à la charge de l'organisation",
                "• PSP/PSA, RBAC, Network Policies à configurer",
                "• Gestion PKI, certificats, secrets (Vault recommandé)",
                "• Audits et conformité sous responsabilité interne",
                "• Outils au choix : Falco, OPA, Trivy, Kyverno…",
                "• Risque élevé si expertise interne insuffisante",
            ],
            "eks": [
                "• IAM AWS + IRSA (IAM Roles for Service Accounts)",
                "• Chiffrement des secrets avec AWS KMS",
                "• AWS GuardDuty pour détection de menaces K8s",
                "• Amazon Inspector pour scan de vulnérabilités",
                "• Pod Security Admission (PSA) activé par défaut",
                "• Conformité SOC 1/2/3, PCI DSS, HIPAA, FedRAMP",
            ],
        },
        {
            "section_num": "03",
            "section_title": "Réseau",
            "dimension": "Réseau & Capacités Ingress",
            "iks": [
                "• CNI : Calico (défaut) — NetworkPolicy supportée",
                "• IBM Cloud Load Balancer (NLB v1/v2) intégré",
                "• Ingress ALB basé sur nginx géré par IBM",
                "• VPC Security Groups + Network ACLs",
                "• Service Mesh : Istio managé (IBM Cloud Service Mesh)",
                "• IBM Cloud Internet Services (CIS) pour CDN/WAF",
            ],
            "roks": [
                "• CNI : OVN-Kubernetes ou SDN (OpenShift Network)",
                "• Routes OpenShift pour exposition des services",
                "• Ingress Kubernetes aussi supporté",
                "• Multus CNI pour interfaces réseau multiples",
                "• Service Mesh : OpenShift Service Mesh (Istio/Kiali/Jaeger)",
                "• HAProxy en tant que contrôleur d'ingress par défaut",
            ],
            "k8s": [
                "• CNI au choix : Calico, Flannel, Cilium, Weave…",
                "• Ingress controller à installer manuellement (nginx, Traefik…)",
                "• Load Balancer externe requis (MetalLB, HAProxy, F5…)",
                "• Flexibilité maximale sur la topologie réseau",
                "• Service Mesh au choix : Istio, Linkerd, Consul…",
                "• Configuration BGP possible avec Calico",
            ],
            "eks": [
                "• CNI : AWS VPC CNI (IPs natives VPC pour pods)",
                "• AWS Load Balancer Controller (ALB/NLB AWS natifs)",
                "• Ingress via ALB Ingress Controller ou nginx",
                "• AWS App Mesh ou Istio pour service mesh",
                "• AWS Transit Gateway pour connectivité multi-VPC",
                "• PrivateLink pour exposition de services privés",
            ],
        },
        {
            "section_num": "03",
            "section_title": "Stockage",
            "dimension": "Options de Stockage",
            "iks": [
                "• IBM Cloud Block Storage (SSD, 10-4000 IOPS)",
                "• IBM Cloud File Storage (NFS v4.1)",
                "• IBM Cloud Object Storage (S3-compatible via S3FS)",
                "• Portworx intégré pour stockage stateful avancé",
                "• Storage Classes pré-configurées",
                "• Snapshots et sauvegardes via IBM Cloud Backup",
            ],
            "roks": [
                "• Même options de stockage IBM Cloud que IKS",
                "• OpenShift Data Foundation (ODF/OCS) disponible",
                "• Portworx supporté nativement",
                "• Dynamic provisioning via StorageClass OCP",
                "• IBM Spectrum Scale possible en option",
                "• Snapshots CSI supportés",
            ],
            "k8s": [
                "• Tout CSI driver compatible supportable",
                "• NFS, iSCSI, FC, Ceph, Longhorn, OpenEBS…",
                "• Stockage objet S3 via Rook-Ceph ou MinIO",
                "• Configuration entièrement manuelle",
                "• Flexibilité totale selon l'infrastructure sous-jacente",
                "• StorageClass à définir par l'équipe",
            ],
            "eks": [
                "• Amazon EBS (gp3, io2 — SSD hautes performances)",
                "• Amazon EFS (NFS managé, multi-AZ)",
                "• Amazon FSx (Lustre, NetApp ONTAP, OpenZFS)",
                "• Amazon S3 via Mountpoint for S3 (CSI natif)",
                "• EBS CSI driver managé par AWS",
                "• Snapshots EBS pour DR",
            ],
        },
        {
            "section_num": "03",
            "section_title": "Coûts",
            "dimension": "Modèle Tarifaire & Coûts",
            "iks": [
                "• Cluster gratuit (1 nœud) pour test",
                "• Facturation par nœud worker (VMs IBM Cloud)",
                "• Plan de contrôle inclus sans surcoût",
                "• Frais réseau sortant (egress IBM Cloud)",
                "• Remises via IBM Cloud commitments",
                "• Coût total maîtrisé pour workloads IBM Cloud",
            ],
            "roks": [
                "• Licence OpenShift incluse dans le tarif IBM Cloud",
                "• Premium vs IKS (~$0.16/h pour licence OCP par nœud)",
                "• Plan de contrôle inclus sans surcoût",
                "• Red Hat Subscription incluse",
                "• Frais réseau sortant IBM Cloud",
                "• Coût supérieur à IKS mais valeur ajoutée OCP",
            ],
            "k8s": [
                "• Coût infrastructure : serveurs, électricité, datacenter",
                "• Coût opérationnel élevé : équipes SRE/Ops dédiées",
                "• Kubernetes open source = 0€ de licence",
                "• ROI si infrastructure existante à forte utilisation",
                "• TCO souvent sous-estimé sans expertise interne",
                "• Peut être économique à très grande échelle",
            ],
            "eks": [
                "• $0.10/h par cluster EKS (plan de contrôle)",
                "• + Coût des instances EC2 (nœuds workers)",
                "• + Fargate : facturation vCPU/mémoire par pod",
                "• Économies via Reserved Instances / Savings Plans",
                "• Frais de transfert de données inter-AZ et sortants",
                "• Coût total potentiellement élevé sans optimisation",
            ],
        },
        {
            "section_num": "03",
            "section_title": "Écosystème",
            "dimension": "Intégrations & Écosystème",
            "iks": [
                "• IBM Cloud Catalog (IA, Db2, Watson, MQ…)",
                "• IBM Cloud Pak (Data, Integration, Security…)",
                "• Terraform IBM provider officiel",
                "• Helm, Kustomize, ArgoCD, Flux supportés",
                "• IBM Cloud Monitoring (Sysdig) intégré",
                "• IBM Cloud Log Analysis (LogDNA) intégré",
            ],
            "roks": [
                "• OperatorHub.io (> 300 opérateurs certifiés Red Hat)",
                "• Red Hat Marketplace (logiciels certifiés OCP)",
                "• IBM Cloud Pak intégrés nativement",
                "• OpenShift GitOps (ArgoCD) et Pipelines (Tekton) natifs",
                "• Intégration Ansible Automation Platform",
                "• Developer Hub / Backstage intégré",
            ],
            "k8s": [
                "• Tout l'écosystème CNCF disponible",
                "• Aucune restriction fournisseur",
                "• Helm, Kustomize, ArgoCD, Flux librement",
                "• Prometheus, Grafana, Loki à installer",
                "• Complexité d'intégration accrue",
                "• Communauté open source mondiale",
            ],
            "eks": [
                "• AWS Marketplace (milliers d'applications K8s)",
                "• Intégration native : ECR, CodePipeline, CodeBuild",
                "• AWS Controllers for Kubernetes (ACK)",
                "• KEDA, Karpenter pour auto-scaling avancé",
                "• AWS Distro for OpenTelemetry (ADOT)",
                "• Ecosystem AWS parmi les plus larges du marché",
            ],
        },
        {
            "section_num": "03",
            "section_title": "Support",
            "dimension": "Support & SLA",
            "iks": [
                "• SLA 99,9% sur le plan de contrôle Kubernetes",
                "• Support IBM Cloud inclus (Basic, Standard, Premium)",
                "• Réponse P1 < 1h (Support Premium)",
                "• Support 24/7 disponible avec abonnement",
                "• Documentation IBM Cloud et communauté IBM",
                "• IBM Client Engineering pour projets stratégiques",
            ],
            "roks": [
                "• SLA 99,9% sur le plan de contrôle OpenShift",
                "• Support IBM Cloud + Red Hat (double couverture)",
                "• Accès au support Red Hat pour OCP",
                "• Réponse P1 < 1h (Support Premium)",
                "• Red Hat Customer Portal + IBM Support",
                "• Support OpenShift inclus dans la licence",
            ],
            "k8s": [
                "• Aucun SLA fournisseur — responsabilité interne",
                "• Support communautaire via CNCF/Slack/GitHub",
                "• Support commercial possible via distribution (Rancher, Tanzu…)",
                "• Temps de résolution dépend des équipes internes",
                "• Nécessite astreinte et runbooks opérationnels",
                "• Formation interne requise",
            ],
            "eks": [
                "• SLA 99,95% sur le plan de contrôle EKS",
                "• Support AWS inclus (Developer, Business, Enterprise)",
                "• AWS Support TAM (Technical Account Manager) disponible",
                "• Réponse P1 < 15 min (Enterprise On-Ramp/Enterprise)",
                "• AWS re:Post, documentation AWS extensive",
                "• Support 24/7 avec plan Business ou supérieur",
            ],
        },
        {
            "section_num": "03",
            "section_title": "Complexité opérationnelle",
            "dimension": "Charge opérationnelle & Complexité de gestion",
            "iks": [
                "• Faible — plan de contrôle géré par IBM",
                "• Gestion nœuds workers (taille, mise à jour)",
                "• Monitoring pré-intégré (Sysdig/LogDNA)",
                "• Auto-scaling intégré (HPA + Cluster Autoscaler)",
                "• Backup etcd automatique par IBM",
                "• Idéal pour équipes de taille moyenne",
            ],
            "roks": [
                "• Faible à Modérée — OCP ajoute une couche d'abstraction",
                "• Gestion Machine Config Operator pour OS nœuds",
                "• Outils OCP supplémentaires à maîtriser (oc, console)",
                "• CVO (Cluster Version Operator) pour mises à jour",
                "• Plus riche fonctionnellement mais plus complexe",
                "• Adapté aux équipes connaissant OpenShift",
            ],
            "k8s": [
                "• Très élevée — toute la pile est gérée en interne",
                "• etcd, API server, scheduler, controller-manager",
                "• Certificats, PKI, CNI, CSI, Ingress à gérer",
                "• Surveillance continue requise (PagerDuty, on-call)",
                "• Nécessite compétences avancées Kubernetes",
                "• Risque opérationnel le plus élevé",
            ],
            "eks": [
                "• Faible — plan de contrôle géré par AWS",
                "• Gestion des Managed Node Groups simplifiée",
                "• Karpenter pour provisionnement intelligent des nœuds",
                "• AWS Fargate pour éliminer la gestion des nœuds",
                "• EKS Add-ons pour composants managés",
                "• Adapté aux équipes AWS expérimentées",
            ],
        },
    ]

    page = 8
    for comp in comparisons:
        _comparison_slide(
            prs,
            section_num=comp["section_num"],
            section_title=comp["section_title"],
            dimension=comp["dimension"],
            iks_text=comp["iks"],
            roks_text=comp["roks"],
            k8s_text=comp["k8s"],
            eks_text=comp["eks"],
            page_num=page
        )
        page += 1

    return page  # retourne le prochain numéro de page


def slide_summary_table(prs, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_slide_header_bar(slide, "04 — Tableau Récapitulatif Comparatif")

    data = [
        ["Critère",                    "IKS",            "ROKS",           "K8s Auto-Géré",     "Amazon EKS"],
        ["Modèle d'hébergement",       "IBM Cloud",      "IBM Cloud",      "On-Prem / IaaS",    "AWS Cloud"],
        ["Facilité de déploiement",    "★★★★☆",          "★★★☆☆",          "★★☆☆☆",             "★★★★☆"],
        ["Gestion des mises à jour",   "★★★★☆",          "★★★★★",          "★★☆☆☆",             "★★★★☆"],
        ["Sécurité managée",           "★★★★☆",          "★★★★★",          "★★★☆☆",             "★★★★☆"],
        ["Écosystème & Intégrations",  "★★★★☆",          "★★★★★",          "★★★★★",             "★★★★★"],
        ["Flexibilité réseau",         "★★★☆☆",          "★★★☆☆",          "★★★★★",             "★★★★☆"],
        ["Options de stockage",        "★★★★☆",          "★★★★★",          "★★★★★",             "★★★★★"],
        ["Coût (maîtrise)",            "★★★★☆",          "★★★☆☆",          "★★★☆☆",             "★★★☆☆"],
        ["Support SLA",                "99,9% – IBM",    "99,9% – IBM+RH", "Interne",           "99,95% – AWS"],
        ["Complexité opérationnelle",  "Faible",         "Modérée",        "Très élevée",       "Faible"],
        ["Conformité (certifications)","ISO/SOC/PCI",    "FIPS/FedRAMP",   "À construire",      "SOC/PCI/HIPAA"],
        ["Vendor lock-in",             "IBM Cloud",      "IBM + Red Hat",  "Minimal",           "AWS"],
        ["Profil idéal",               "IBM Cloud natif","Enterprise RH",  "Souveraineté data", "AWS natif"],
    ]

    add_table(slide, data, x=0.15, y=1.2, w=13.0, h=6.0,
              header_fill=DARK_BLUE, header_text=WHITE,
              row_alt=LIGHT_BLUE, font_size=11)

    add_footer(slide, page_num=str(page_num))


def slide_recommendations(prs, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_slide_header_bar(slide, "05 — Recommandations par Profil Organisationnel")

    profiles = [
        {
            "title": "Entreprise sur IBM Cloud",
            "color": MED_BLUE,
            "solution": "→ IKS ou ROKS",
            "bullets": [
                "Utilise déjà IBM Cloud (Watson, Db2, Cloud Paks…)",
                "ROKS si équipes Red Hat et CI/CD natif requis",
                "IKS si Kubernetes standard suffisant",
                "SLA IBM + support intégré",
            ]
        },
        {
            "title": "Cloud Hybride / Multi-Cloud",
            "color": RED_OCP,
            "solution": "→ ROKS (OpenShift) ou K8s Auto-Géré",
            "bullets": [
                "ROKS avec OpenShift : portabilité on-prem ↔ cloud",
                "ACM (Advanced Cluster Management) pour multi-clusters",
                "K8s auto-géré si contrôle total requis on-prem",
                "Cohérence oc / OCP entre environnements",
            ]
        },
        {
            "title": "Organisation AWS-Native",
            "color": ORANGE_EKS,
            "solution": "→ Amazon EKS",
            "bullets": [
                "Déjà investi dans l'écosystème AWS (IAM, VPC, RDS…)",
                "Fargate pour workloads serverless sans gestion nœuds",
                "Karpenter pour optimisation des coûts EC2",
                "Intégration native CodePipeline, CodeBuild, ECR",
            ]
        },
        {
            "title": "Souveraineté des données / On-Premises",
            "color": GREEN_K8S,
            "solution": "→ K8s Auto-Géré",
            "bullets": [
                "Exigences réglementaires interdisant le cloud public",
                "Infrastructure propre existante à valoriser",
                "Envisager Rancher, VMware Tanzu ou OpenShift autogéré",
                "Investissement fort en expertise SRE/Ops requis",
            ]
        },
    ]

    for i, prof in enumerate(profiles):
        col = i % 2
        row = i // 2
        bx = 0.25 + col * 6.55
        by = 1.2 + row * 3.05

        add_rect(slide, bx, by, 6.3, 0.45, fill_color=prof["color"])
        add_text_box(slide, prof["title"], bx, by, 6.3, 0.45,
                     font_size=15, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, font_name="Calibri")

        add_rect(slide, bx, by + 0.45, 6.3, 2.55, fill_color=LIGHT_GRAY)
        add_text_box(slide, prof["solution"], bx + 0.1, by + 0.5, 6.0, 0.4,
                     font_size=14, bold=True, color=prof["color"], font_name="Calibri")

        add_colored_bullet(slide, prof["bullets"],
                           bx + 0.1, by + 0.95, 6.0, 2.0,
                           font_size=13, bullet_color=prof["color"], font_name="Calibri")

    add_footer(slide, page_num=str(page_num))


def slide_conclusion(prs, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BLUE)

    # Titre
    add_rect(slide, 0, 0, 13.33, 1.2, fill_color=MED_BLUE)
    add_text_box(slide, "06 — Conclusion & Points Clés",
                 0.4, 0.15, 12.5, 0.9,
                 font_size=32, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT, font_name="Calibri")

    takeaways = [
        ("Aucune solution universelle",
         "Le choix dépend du contexte : cloud existant, exigences réglementaires, expertise interne et budget."),
        ("IKS : simplicité sur IBM Cloud",
         "Idéal pour les organisations IBM Cloud cherchant Kubernetes standard avec un faible TCO opérationnel."),
        ("ROKS : richesse entreprise",
         "La meilleure option pour les environnements hybrides Red Hat, avec sécurité et CI/CD intégrés."),
        ("K8s Auto-Géré : contrôle maximal",
         "Recommandé uniquement si la souveraineté des données ou l'absence de cloud public est obligatoire."),
        ("Amazon EKS : leader AWS",
         "Le choix naturel pour les organisations AWS-native, avec un écosystème et une maturité incomparables."),
        ("Coût total de propriété (TCO)",
         "Ne pas négliger le coût opérationnel : K8s auto-géré peut coûter plus cher que les solutions managées."),
    ]

    for i, (title, body) in enumerate(takeaways):
        row = i // 2
        col = i % 2
        bx = 0.3 + col * 6.5
        by = 1.35 + row * 1.95

        add_rect(slide, bx, by, 6.3, 0.45, fill_color=ACCENT_GOLD)
        add_text_box(slide, title, bx + 0.1, by + 0.03, 6.1, 0.4,
                     font_size=14, bold=True, color=DARK_BLUE, font_name="Calibri")
        add_rect(slide, bx, by + 0.45, 6.3, 1.45, fill_color=RGBColor(0x16, 0x1D, 0x52))
        add_text_box(slide, body, bx + 0.15, by + 0.5, 6.0, 1.35,
                     font_size=13, bold=False, color=WHITE, font_name="Calibri")

    add_footer(slide, page_num=str(page_num))


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_01_title(prs)
    slide_02_agenda(prs)
    slide_03_intro(prs)
    slide_04_iks(prs)
    slide_05_roks(prs)
    slide_06_k8s_self(prs)
    slide_07_eks(prs)
    next_page = slides_comparative(prs)        # slides 8 → 17
    slide_summary_table(prs, page_num=next_page)
    slide_recommendations(prs, page_num=next_page + 1)
    slide_conclusion(prs, page_num=next_page + 2)

    output = "IKS-vs-ROKS-Analyse-Comparative.pptx"
    prs.save(output)
    print(f"✅ Présentation sauvegardée : {output}")
    print(f"   Nombre de diapositives : {len(prs.slides)}")


if __name__ == "__main__":
    main()
